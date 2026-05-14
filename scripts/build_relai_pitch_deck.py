from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "presentation"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUT_DIR / "RelAi_pitch_deck_draft.pptx"

IMG_CONCEPT = ROOT / "docs" / "presentation" / "assets" / "relai-brand-concept-board.png"
IMG_STORY = ROOT / "docs" / "presentation" / "assets" / "relai-ppt-storyboard-concept.png"
IMG_STEP1 = ROOT / "Relai" / "step1.png"
IMG_STEP5 = ROOT / "Relai" / "step5.png"
IMG_STEP68 = ROOT / "Relai" / "step6,8.png"
IMG_CONTROL = ROOT / "relai-mvp-control-tower.png"


BG = RGBColor(255, 255, 255)
TEXT = RGBColor(16, 16, 16)
SUBTLE = RGBColor(102, 102, 102)
LINE = RGBColor(220, 224, 230)
BLUE = RGBColor(78, 105, 255)
PALE_BLUE = RGBColor(243, 246, 255)
BLACK = RGBColor(0, 0, 0)

W = Inches(13.333)
H = Inches(7.5)


def add_textbox(slide, x, y, w, h, text, *, size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if not text:
        return box
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_body_lines(slide, x, y, w, lines, *, size=18, color=TEXT, bullet=False, gap=0.28):
    box = slide.shapes.add_textbox(x, y, w, Inches(3.8))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {line}" if bullet else line
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(6.9), Inches(0.5), title, size=28, bold=True)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.02), Inches(1.45), Inches(0.05)
    ).fill.solid()
    line = slide.shapes[-1]
    line.line.fill.background()
    line.fill.fore_color.rgb = BLUE
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.15), Inches(8.8), Inches(0.45), subtitle, size=12, color=SUBTLE)


def card(slide, x, y, w, h, title, body="", *, fill=BG, title_size=14, body_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.1), w - Inches(0.24), Inches(0.22), title, size=title_size, bold=True)
    if body:
        add_textbox(
            slide,
            x + Inches(0.12),
            y + Inches(0.38),
            w - Inches(0.24),
            h - Inches(0.44),
            body,
            size=body_size,
            color=SUBTLE,
        )
    return shape


def stat_card(slide, x, y, w, h, value, label, note):
    shape = card(slide, x, y, w, h, label, "", fill=PALE_BLUE, title_size=12)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), w - Inches(0.3), Inches(0.45), value, size=24, bold=True, color=BLUE)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.87), w - Inches(0.3), Inches(0.4), note, size=10, color=SUBTLE)
    return shape


def arrow_label(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.color.rgb = LINE
    add_textbox(slide, x + Inches(0.08), y + Inches(0.11), w - Inches(0.16), Inches(0.2), text, size=11, bold=True, align=PP_ALIGN.CENTER)


def image(slide, path: Path, x, y, w=None, h=None):
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def footer(slide, text):
    add_textbox(slide, Inches(0.75), Inches(7.0), Inches(11.9), Inches(0.2), text, size=9, color=SUBTLE)


def make_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    add_textbox(s, Inches(0.8), Inches(0.9), Inches(5.5), Inches(0.6), "RelAi", size=30, bold=True)
    add_textbox(
        s,
        Inches(0.8),
        Inches(1.6),
        Inches(5.8),
        Inches(1.0),
        "요구사항 변경을 팀의 실행 단위로 정렬하는\nAI 협업 운영 시스템",
        size=22,
        bold=True,
    )
    add_body_lines(s, Inches(0.85), Inches(3.0), Inches(4.9), ["Alignment", "Memory", "Impact", "Approval", "Orchestration"], size=16)
    image(s, IMG_CONCEPT, Inches(6.45), Inches(0.7), w=Inches(6.0))

    # 2 Problem
    s = prs.slides.add_slide(blank)
    add_title(s, "문제는 변경 그 자체가 아니라, 합의 구조의 부재입니다")
    roles = ["기획자", "개발자", "QA", "승인권자"]
    tags = ["요구 변경", "재작업", "검수 수정", "승인 지연"]
    for i, role in enumerate(roles):
        card(s, Inches(0.9 + i * 3.05), Inches(2.0), Inches(2.45), Inches(1.05), role, "서로 다른 관점으로\n같은 요구사항을 해석")
    for i, tag in enumerate(tags[:-1]):
        arrow_label(s, Inches(2.95 + i * 3.05), Inches(2.35), Inches(0.8), Inches(0.35), tag)
    add_body_lines(
        s,
        Inches(0.95),
        Inches(3.55),
        Inches(11.4),
        [
            "기획 초안이 나오면 개발이 시작되지만, 개발 중간에 기획이 수정되고 품질 검수 과정에서 다시 수정이 열립니다.",
            "문제는 누군가의 역량 부족이 아니라, 변경을 같은 기준으로 다시 정렬하는 구조가 없다는 데 있습니다.",
        ],
        size=18,
    )
    footer(s, "현업 병목: 변경 발생 -> 영향 파악 지연 -> 회의 반복 -> 승인 지연 -> 재작업 증가")

    # 3 Why now
    s = prs.slides.add_slide(blank)
    add_title(s, "AI는 빨라졌지만, 팀의 정렬 방식은 그대로입니다", "개인 생산성 상승과 조직 정렬 정체의 간극")
    stat_card(s, Inches(0.9), Inches(1.8), Inches(2.7), Inches(1.45), "75%", "지식근로자 AI 사용", "Microsoft 2024")
    stat_card(s, Inches(3.8), Inches(1.8), Inches(2.7), Inches(1.45), "68%", "개발자 시간 절약", "Atlassian 2025")
    stat_card(s, Inches(6.7), Inches(1.8), Inches(2.7), Inches(1.45), "80%", "메시지 오해로 재작업", "Atlassian/Loom 2026")
    stat_card(s, Inches(9.6), Inches(1.8), Inches(2.7), Inches(1.45), "40%", "요구사항 수집 부실", "PMI 2021")
    card(s, Inches(1.0), Inches(3.7), Inches(5.6), Inches(2.0), "개인 단위", "AI로 문서, 코드, 테스트를 더 빨리 만든다", fill=PALE_BLUE)
    card(s, Inches(6.75), Inches(3.7), Inches(5.6), Inches(2.0), "조직 단위", "팀은 여전히 문서, 메신저, 회의 중심으로 합의를 맞춘다", fill=BG)
    arrow_label(s, Inches(5.88), Inches(4.45), Inches(0.7), Inches(0.36), "간극")

    # 4 Insight
    s = prs.slides.add_slide(blank)
    add_title(s, "진짜 AX는 개인이 아니라 조직의 의사소통 속도를 높이는 것입니다")
    card(s, Inches(1.0), Inches(1.85), Inches(11.2), Inches(3.4), "", "", fill=PALE_BLUE)
    add_textbox(
        s,
        Inches(1.35),
        Inches(2.35),
        Inches(10.5),
        Inches(1.2),
        "개인의 AI 활용이 팀의 합의 구조로 연결되지 않으면,\n생산성은 빨라져도 협업 비용은 그대로 남습니다.",
        size=24,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.8),
        Inches(4.1),
        Inches(9.8),
        Inches(0.8),
        "RelAi는 요구사항 변경, 팀 논의, 승인 과정을 같은 기준으로 다시 정렬하는 시스템입니다.",
        size=18,
        color=SUBTLE,
        align=PP_ALIGN.CENTER,
    )

    # 5 Solution
    s = prs.slides.add_slide(blank)
    add_title(s, "RelAi는 변경을 실행 가능한 합의 구조로 바꿉니다")
    card(s, Inches(0.9), Inches(1.8), Inches(2.9), Inches(2.1), "입력", "요구사항 문서\n회의 기록\n팀 채팅\n변경 요청", fill=BG, title_size=18, body_size=16)
    arrow_label(s, Inches(4.0), Inches(2.5), Inches(1.0), Inches(0.4), "AI 구조화")
    card(s, Inches(5.2), Inches(1.8), Inches(3.0), Inches(2.1), "처리", "역할별 해석\n합의 상태 분리\n영향 분석\n메모리 연결", fill=PALE_BLUE, title_size=18, body_size=16)
    arrow_label(s, Inches(8.45), Inches(2.5), Inches(1.0), Inches(0.4), "출력")
    card(s, Inches(9.7), Inches(1.8), Inches(2.7), Inches(2.1), "결과", "결정\n질문\n승인\n영향\n액션\n메모리", fill=BG, title_size=18, body_size=16)
    add_body_lines(s, Inches(0.95), Inches(4.4), Inches(11.2), ["문서를 더 많이 만드는 것이 아니라, 이미 존재하는 논의를 팀이 실제로 움직일 수 있는 구조로 바꿉니다."], size=18)

    # 6 Core flow
    s = prs.slides.add_slide(blank)
    add_title(s, "초안부터 변경, 영향, 회의까지 하나의 흐름으로 이어집니다")
    steps = ["초안 등록", "역할별 해석", "합의 잠금", "영향 맵", "프로젝트 메모리", "AI 회의 개설"]
    x = 0.75
    for step in steps:
        card(s, Inches(x), Inches(1.7), Inches(1.95), Inches(1.05), step, "", fill=PALE_BLUE if step in {"영향 맵", "프로젝트 메모리"} else BG, title_size=15)
        x += 2.1
    for i in range(5):
        arrow_label(s, Inches(2.52 + i * 2.1), Inches(2.07), Inches(0.42), Inches(0.3), "")
    image(s, IMG_STORY, Inches(1.1), Inches(3.1), w=Inches(11.1))
    footer(s, "핵심 플로우: 입력 -> 해석 -> 분리 -> 영향 -> 누적 -> 오케스트레이션")

    # 7 Scenario
    s = prs.slides.add_slide(blank)
    add_title(s, "변경 요청 하나가 어떻게 정렬되는가")
    card(s, Inches(0.85), Inches(1.7), Inches(3.2), Inches(1.1), "변경 요청", "관리자 승인 후 결제 확정으로 정책 변경", fill=PALE_BLUE, title_size=17, body_size=14)
    image(s, IMG_STEP5, Inches(4.35), Inches(1.6), w=Inches(7.2))
    add_body_lines(
        s,
        Inches(0.95),
        Inches(3.2),
        Inches(3.0),
        [
            "영향 화면: 결제 완료 / 관리자 승인",
            "영향 API: payment-confirm",
            "영향 데이터: approval_status",
            "테스트 영향: 승인 전/후 QA",
            "일정 영향: QA 2일 증가 가능",
        ],
        size=14,
        bullet=True,
    )

    # 8 Product screens
    s = prs.slides.add_slide(blank)
    add_title(s, "실제 MVP 화면에서 정렬의 흐름이 보입니다")
    image(s, IMG_CONTROL, Inches(0.85), Inches(1.55), w=Inches(6.2))
    image(s, IMG_STEP1, Inches(7.3), Inches(1.55), w=Inches(5.15))
    image(s, IMG_STEP68, Inches(7.3), Inches(4.15), w=Inches(5.15))
    add_textbox(s, Inches(0.9), Inches(6.35), Inches(11.2), Inches(0.35), "대시보드 / 영향 맵 / 프로젝트 메모리 / AI 회의 개설까지 한 흐름으로 연결", size=13, color=SUBTLE)

    # 9 Differentiation
    s = prs.slides.add_slide(blank)
    add_title(s, "기존 도구는 기록을 남기고, RelAi는 다음 실행을 정렬합니다")
    headers = ["구분", "기존 도구", "RelAi"]
    rows = [
        ("Mantis / Jira", "이슈를 남김", "변경이 누구에게 어떤 의미인지 해석"),
        ("Confluence / Notion", "문서를 남김", "문서와 대화를 프로젝트 메모리로 연결"),
        ("메신저 / 회의", "대화를 남김", "결정, 질문, 승인, 액션으로 구조화"),
    ]
    col_x = [0.95, 3.0, 6.35]
    col_w = [1.8, 3.0, 5.15]
    for i, head in enumerate(headers):
        card(s, Inches(col_x[i]), Inches(1.75), Inches(col_w[i]), Inches(0.55), head, "", fill=PALE_BLUE, title_size=16)
    y = 2.45
    for row in rows:
        for i, value in enumerate(row):
            card(s, Inches(col_x[i]), Inches(y), Inches(col_w[i]), Inches(0.8), value, "", fill=BG, title_size=13)
        y += 0.88
    add_body_lines(s, Inches(0.95), Inches(5.45), Inches(11.0), ["기존 도구가 기록을 축적하는 역할을 한다면, RelAi는 그 기록을 바탕으로 다음에 누가 무엇을 결정하고 실행해야 하는지까지 정렬합니다."], size=17)

    # 10 Feasibility and business
    s = prs.slides.add_slide(blank)
    add_title(s, "RelAi는 이미 검증 가능한 단계에 와 있습니다")
    card(s, Inches(0.9), Inches(1.7), Inches(3.65), Inches(1.85), "실현 가능성", "웹 기반 MVP 개발 중\n핵심 기능 완성 단계\n실제 업무 흐름 기준 검증", fill=PALE_BLUE, title_size=18, body_size=15)
    card(s, Inches(4.85), Inches(1.7), Inches(3.65), Inches(1.85), "초기 고객", "제품팀\n개발팀\nPM 조직\nSI/외주 협업팀", fill=BG, title_size=18, body_size=15)
    card(s, Inches(8.8), Inches(1.7), Inches(3.6), Inches(1.85), "수익 모델", "팀 구독\n프로젝트 과금\n엔터프라이즈\n부서 확장", fill=BG, title_size=18, body_size=15)
    add_body_lines(
        s,
        Inches(0.95),
        Inches(4.2),
        Inches(11.2),
        [
            "초기에는 독립형 시스템으로 검증하고, 이후 Jira, Notion, Slack, Teams, GitHub, Mantis, Confluence와의 연동으로 실제 조직 안으로 자연스럽게 들어갈 수 있습니다.",
            "핵심 가치는 문서 작성이 아니라 재작업 감소, 승인 속도 개선, 의사결정 메모리 축적입니다.",
        ],
        size=17,
    )

    # 11 Closing
    s = prs.slides.add_slide(blank)
    add_title(s, "더 빠른 AI 시대에는, 더 빠른 합의 구조가 필요합니다")
    card(s, Inches(1.05), Inches(1.7), Inches(11.0), Inches(3.4), "", "", fill=PALE_BLUE)
    add_textbox(
        s,
        Inches(1.55),
        Inches(2.35),
        Inches(10.0),
        Inches(1.4),
        "RelAi는 더 빠른 변경을,\n더 빠른 합의로 바꾸는 시스템입니다.",
        size=28,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.8),
        Inches(4.25),
        Inches(9.6),
        Inches(0.8),
        "개인의 AI 활용을 넘어, 조직의 의사소통 속도까지 AI 흐름에 맞게 끌어올리는 것이 RelAi의 목표입니다.",
        size=17,
        color=SUBTLE,
        align=PP_ALIGN.CENTER,
    )

    # Appendix A
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix A. 문제 근거 수치")
    items = [
        ("Microsoft 2024", "지식근로자 75% AI 사용, AI 사용자 78% BYOAI"),
        ("Gallup 2025", "업무상 AI 사용 2023년 21% -> 2025년 40%"),
        ("Atlassian 2025", "개발자 68% AI로 시간 절약, 50%는 협업 비효율로 시간 손실"),
        ("Atlassian/Loom 2026", "근로자 80%가 메시지 오해로 재작업 경험"),
        ("PMI 2021", "요구사항 수집 부실 40%, 목표 변경 38%, 커뮤니케이션 부실 28%"),
    ]
    y = 1.7
    for title, body in items:
        card(s, Inches(0.9), Inches(y), Inches(11.5), Inches(0.72), title, body, fill=BG, title_size=14, body_size=12)
        y += 0.82

    # Appendix B
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix B. 제품 구조 상세")
    card(s, Inches(0.9), Inches(1.75), Inches(3.6), Inches(2.4), "입력 레이어", "문서\n채팅\n회의록\n변경 요청", fill=BG, title_size=18, body_size=16)
    card(s, Inches(4.85), Inches(1.75), Inches(3.6), Inches(2.4), "처리 레이어", "역할 변환\n영향 분석\n메모리화\n회의 오케스트레이션", fill=PALE_BLUE, title_size=18, body_size=16)
    card(s, Inches(8.8), Inches(1.75), Inches(3.55), Inches(2.4), "출력 레이어", "결정\n질문\n승인\n영향\n액션", fill=BG, title_size=18, body_size=16)
    add_body_lines(s, Inches(0.95), Inches(4.6), Inches(11.0), ["RelAi의 핵심은 요약이 아니라 구조화다. 입력된 협업 맥락을 다음 실행이 가능한 단위로 바꾼다."], size=17)

    # Appendix C
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix C. 프로젝트 메모리")
    card(s, Inches(0.95), Inches(1.75), Inches(5.2), Inches(3.5), "저장되는 것", "회의 기록\n채팅\n기술 제약\n승인 이력\n반복 쟁점", fill=BG, title_size=20, body_size=18)
    card(s, Inches(6.45), Inches(1.75), Inches(5.0), Inches(3.5), "활용되는 방식", "새 변경 요청이 들어오면\n과거 합의와 충돌 여부 판단\n관련 근거 자동 연결\n다음 질문과 승인 항목 제안", fill=PALE_BLUE, title_size=20, body_size=18)

    # Appendix D
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix D. 기존 도구와의 관계")
    add_body_lines(
        s,
        Inches(0.95),
        Inches(1.8),
        Inches(11.1),
        [
            "초기에는 독립형 시스템으로 시작할 수 있습니다.",
            "장기적으로는 Jira, Notion, Slack, Teams, GitHub, Mantis, Confluence와 연동합니다.",
            "기존 도구를 버리게 하는 것이 아니라, 그 위에서 AI 정렬 레이어로 작동합니다.",
        ],
        size=20,
        bullet=True,
    )

    # Appendix E
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix E. 시장 진입 순서")
    tiers = [
        ("1차", "제품팀 / 개발팀 / PM 조직"),
        ("2차", "SI / 외주 협업 프로젝트"),
        ("3차", "사내 디지털 전환 조직"),
        ("4차", "마케팅 / 운영 / 고객지원 등 타 부서"),
    ]
    y = 1.75
    for idx, (a, b) in enumerate(tiers):
        card(s, Inches(1.0), Inches(y + idx * 1.05), Inches(1.15), Inches(0.7), a, "", fill=PALE_BLUE, title_size=16)
        card(s, Inches(2.35), Inches(y + idx * 1.05), Inches(9.1), Inches(0.7), b, "", fill=BG, title_size=16)

    # Appendix F
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix F. 리스크와 대응")
    rows = [
        ("AI가 잘못 해석할 수 있음", "사람 승인 구조 유지, 잠긴 결정과 열린 질문 분리"),
        ("근거 없는 추천 위험", "프로젝트 메모리와 근거 연결 구조 사용"),
        ("도입 장벽", "독립형 MVP로 시작 후 기존 툴 연동 확장"),
    ]
    y = 1.8
    for risk, action in rows:
        card(s, Inches(0.95), Inches(y), Inches(4.2), Inches(0.95), risk, "", fill=BG, title_size=15)
        card(s, Inches(5.45), Inches(y), Inches(6.9), Inches(0.95), action, "", fill=PALE_BLUE, title_size=15)
        y += 1.12

    # Appendix G
    s = prs.slides.add_slide(blank)
    add_title(s, "Appendix G. 예상 질문")
    qas = [
        ("왜 기존 협업툴로 안 되나?", "기록은 남기지만, 변경 발생 시 다음 실행까지 정렬해주지 못합니다."),
        ("AI가 오판하면 더 위험하지 않나?", "RelAi는 자동 실행이 아니라 판단 지원이며, 사람 승인 구조를 유지합니다."),
        ("가장 먼저 누가 돈 내고 쓸까?", "제품팀, 개발팀, PM 조직, SI/외주 협업팀이 가장 적합합니다."),
        ("장기적으로 어떤 서비스가 되나?", "조직의 요구사항 변경과 의사결정 메모리를 관리하는 AX 협업 플랫폼으로 확장됩니다."),
    ]
    y = 1.55
    for q, a in qas:
        card(s, Inches(0.9), Inches(y), Inches(11.5), Inches(1.02), q, a, fill=BG, title_size=14, body_size=11)
        y += 1.12

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    make_presentation()
